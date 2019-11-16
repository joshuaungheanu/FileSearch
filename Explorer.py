import sys, os, re
import openpyxl
from xlsxwriter.utility import xl_rowcol_to_cell
import xlrd
import docx
from docx import Document
from pptx import Presentation
from PyQt5 import QtCore, QtGui
from PyQt5.QtWidgets import *
from PyQt5.QtGui import QTextCursor, QTextCharFormat, QBrush, QColor
from PyQt5.QtCore import pyqtSlot


import re

class App(QMainWindow):

    def __init__(self):
        super().__init__()
        self.title = 'File Explorer'
        self.left = 10
        self.top = 10
        self.width = 680
        self.height = 660
        self.all_drives = self.get_drives()
        self.initUI()

    def initUI(self):
        self.all_drives = self.get_drives()

        self.setWindowTitle(self.title)
        self.setGeometry(self.left, self.top, self.width, self.height)

        self.button = QPushButton('Search', self)
        self.button.move(550, 25)

        self.combo = QComboBox(self)
        self.combo.resize(500, 25)
        self.combo.move(20, 80)

        self.button2 = QPushButton('Go', self)
        self.button2.move(550, 77)

        # Create textbox
        self.textbox = QLineEdit(self)
        self.textbox.move(20, 20)
        self.textbox.resize(500, 40)

        self.response_box = QTextEdit(self)
        self.response_box.setReadOnly(True)
        self.response_box.move(20, 120)
        self.response_box.resize(640, 500)

        self.textbox.returnPressed.connect(self.button.click)
        self.combo.activated.connect(self.button2.click)
        self.button2.clicked.connect(self.get_text)
        self.button.clicked.connect(self.on_click)

        self.show()

    @pyqtSlot()
    def get_text(self):
        self.response_box.setText('')
        chosen_path = self.combo.currentText()
        if chosen_path.endswith('.txt'):
            text_content = list()
            key_content = ""
            non_key_content = ""
            with open(chosen_path) as f:
                text = f.readlines()
                for line in text:
                    fields = line.split(" ")
                    for word in fields:
                        if word != self.keyword:
                            text_content.append(word)
                        else:
                            non_key_content = " ".join(text_content)
                            self.response_box.setTextColor(QColor(0, 0, 0))
                            self.response_box.insertPlainText(non_key_content)
                            del text_content[:]
                            key_content = " " + self.keyword
                            non_key_content = ""
                            self.response_box.setTextColor(QColor(255, 0, 0))
                            self.response_box.insertPlainText(key_content)
                    if text_content:
                        non_key_content = " " + " ".join(text_content)
                        self.response_box.setTextColor(QColor(0, 0, 0))
                        self.response_box.insertPlainText(non_key_content)
                        del text_content[:]
                    text_content.append("\n")

        elif chosen_path.endswith('.xlsx'):
            normalData = list()
            normalOutput = ""
            key_content = ""
            wb = xlrd.open_workbook(chosen_path)
            sheet = wb.sheet_by_index(0)
            for row_num in range(sheet.nrows):
                for col_num in range(sheet.ncols):
                    cell_obj = sheet.row(row_num)[col_num]
                    if cell_obj.value != self.keyword:
                        normalData.append(cell_obj.value)
                    else:
                        normalOutput = "    ".join(normalData)
                        self.response_box.setTextColor(QColor(0, 0, 0))
                        self.response_box.insertPlainText(normalOutput)
                        del normalData[:]
                        key_content = self.keyword + "  "
                        self.response_box.setTextColor(QColor(255, 0, 0))
                        self.response_box.insertPlainText(key_content)
                normalData.append('\n')
            if normalData:
                normalOutput = "    ".join(normalData)
                self.response_box.setTextColor(QColor(0, 0, 0))
                self.response_box.insertPlainText(normalOutput)
        elif chosen_path.endswith('.pptx'):
            normalData = list()
            normalOutput = ""
            key_content = ""
            tmp_text = ""
            prs = Presentation(chosen_path)
            for slides in prs.slides:
                for shape in slides.shapes:
                    if shape.has_text_frame:
                        if (shape.text.find(self.keyword)) != -1:
                            tmp_text += shape.text
                            fields = tmp_text.split(" ")
                            for word in fields:
                                if word != self.keyword:
                                    normalData.append(word)
                                else:
                                    normalOutput = " ".join(normalData)
                                    normalOutput += " "
                                    self.response_box.setTextColor(QColor(0, 0, 0))
                                    self.response_box.insertPlainText(normalOutput)
                                    del normalData[:]
                                    key_content = " " + self.keyword + " "
                                    normalOutput = ""
                                    self.response_box.setTextColor(QColor(255, 0, 0))
                                    self.response_box.insertPlainText(key_content)
                            if normalData:
                                normalOutput = " ".join(normalData)
                                self.response_box.setTextColor(QColor(0, 0, 0))
                                self.response_box.insertPlainText(normalOutput)
                                normalOutput = ""
                        else:
                            normalOutput += shape.text
                            self.response_box.setTextColor(QColor(0, 0, 0))
                            self.response_box.insertPlainText(normalOutput)
                            normalOutput = ""
                        tmp_text = ""
                        normalOutput += "\n"
        elif chosen_path.endswith('.docx'):
            normalData = list()
            normalOutput = ""
            key_content = ""
            tmp_text = ""
            f = chosen_path
            document = Document(f)
            for p in document.paragraphs:
                if p.text.find(self.keyword) != -1:
                    tmp_text += p.text
                    fields = tmp_text.split(" ")
                    for word in fields:
                        if word != self.keyword:
                            normalData.append(word)
                        else:
                            normalOutput = " ".join(normalData)
                            normalOutput += " "
                            self.response_box.setTextColor(QColor(0, 0, 0))
                            self.response_box.insertPlainText(normalOutput)
                            del normalData[:]
                            key_content = " " + self.keyword + " "
                            normalOutput = ""
                            self.response_box.setTextColor(QColor(255, 0, 0))
                            self.response_box.insertPlainText(key_content)
                    if normalData:
                        normalData.append("\n")
                        normalOutput = " " + " ".join(normalData)
                        self.response_box.setTextColor(QColor(0, 0, 0))
                        self.response_box.insertPlainText(normalOutput)
                        del normalData[:]
                        normalOutput = ""
                        tmp_text = ""
                    else:
                        normalData.append("\n")
                else:
                    normalOutput += p.text
                    normalOutput += "\n"
                    self.response_box.setTextColor(QColor(0, 0, 0))
                    self.response_box.insertPlainText(normalOutput)
                    normalOutput = ""




        self.combo.clear()

    @pyqtSlot()
    def on_click(self):
        textboxValue = self.textbox.text()
        self.keyword = textboxValue
        li = self.search_directory(textboxValue)
        self.combo.addItems(li)
        self.textbox.setText("")

    @staticmethod
    def get_drives():
        response = os.popen("wmic logicaldisk get caption")
        list1 = []
        for line in response.readlines():
            line = line.strip("\n")
            line = line.strip("\r")
            line = line.strip(" ")
            if (line == "Caption" or line == ""):
                continue
            list1.append(line + '/')
        return list1

    def search_directory(self, keyword):
        results = []
        for each in self.all_drives:
            for root, dir, files in os.walk(each, topdown=True):
                for f in files:
                    if os.path.splitext(f)[1] == '.txt':
                        cur_f = open(os.path.join(root,f))
                        #cur_f = open(os.path.join(root, f))
                        if cur_f.read().find(keyword):
                            results.append(os.path.join(root, f))

                    if os.path.splitext(f)[1] == '.xlsx':
                        #wb = xlrd.open_workbook(os.path.expanduser('C:/a1/a1.xlsx'))
                        wb = xlrd.open_workbook(os.path.join(root,f))
                        sheet = wb.sheet_by_index(0)
                        for row_num in range(sheet.nrows):
                            for col_num in range(sheet.ncols):
                                cell_obj = sheet.row(row_num)[col_num]
                                if keyword == cell_obj.value:
                                    results.append(os.path.join(root, f))
                    if os.path.splitext(f)[1] == '.pptx':
                        #f = open(os.path.expanduser('~/.' + f)
                        prs = Presentation(os.path.join(root,f))
                        for slides in prs.slides:
                            for shape in slides.shapes:
                                if shape.has_text_frame:
                                    if (shape.text.find(keyword)) != -1:
                                        results.append(os.path.join(root, f))
                    if os.path.splitext(f)[1] == '.docx':
                        f = os.path.join(root,f)
                        #f = os.path.join(root, f)
                        document = Document(f)
                        for p in document.paragraphs:
                            if p.text.find(keyword) != -1:
                                results.append(os.path.join(root, f))
        return results


if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = App()
    sys.exit(app.exec_())
